from flask import Flask, request, render_template, send_file, after_this_request
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
import os
from fpdf import FPDF
from PIL import Image
import pandas as pd
import docx
import markdown2
import traceback
import json
import xml.etree.ElementTree as ET
from weasyprint import HTML
from pptx import Presentation
import tempfile
import threading
import time
import requests  # ✅ For reCAPTCHA

UPLOAD_FOLDER = os.path.join(os.getcwd(), "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10 MB limit

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    try:
        # ✅ Verify reCAPTCHA
        recaptcha_token = request.form.get('recaptcha_token')
        recaptcha_secret = "6LdU8IciAAAAAFm5f4o9qu_3NWe3z3_XN2sNss41"
        recaptcha_url = "https://www.google.com/recaptcha/api/siteverify"
        recaptcha_response = requests.post(
            recaptcha_url,
            data={'secret': recaptcha_secret, 'response': recaptcha_token}
        )
        result = recaptcha_response.json()
        if not result.get("success") or result.get("score", 0) < 0.5:
            return "reCAPTCHA verification failed. Please try again.", 403

        # ✅ Proceed with file processing
        file = request.files.get('file')
        if not file or file.filename == '':
            return "No file selected"

        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)

        ext = os.path.splitext(filename)[1].lower()
        output_pdf = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.pdf")

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)

        if ext == ".txt":
            with open(filepath, 'r', encoding='utf-8') as f:
                for line in f:
                    pdf.cell(200, 10, txt=line.strip(), ln=1)
            pdf.output(output_pdf)

        elif ext == ".csv":
            df = pd.read_csv(filepath)
            for line in df.to_string(index=False).split('\n'):
                pdf.cell(200, 8, txt=line.strip(), ln=1)
            pdf.output(output_pdf)

        elif ext == ".xlsx":
            df = pd.read_excel(filepath)
            for line in df.to_string(index=False).split('\n'):
                pdf.cell(200, 8, txt=line.strip(), ln=1)
            pdf.output(output_pdf)

        elif ext == ".docx":
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                pdf.cell(200, 10, txt=para.text.strip(), ln=1)
            pdf.output(output_pdf)

        elif ext in [".jpg", ".jpeg", ".png"]:
            img = Image.open(filepath)
            img_rgb = img.convert('RGB')
            img_rgb.save(output_pdf)

        elif ext == ".pptx":
            try:
                prs = Presentation(filepath)
                for slide in prs.slides:
                    text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    for line in text.split('\n'):
                        pdf.cell(200, 10, txt=line.strip(), ln=1)
                pdf.output(output_pdf)
            except Exception as e:
                return f"Error reading .pptx file: {e}"

        elif ext == ".html":
            HTML(filename=filepath).write_pdf(output_pdf)

        elif ext == ".md":
            with open(filepath, 'r', encoding='utf-8') as f:
                html = markdown2.markdown(f.read())
            with tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode='w', encoding='utf-8') as tmp:
                tmp.write(html)
                tmp.close()
                HTML(tmp.name).write_pdf(output_pdf)

        elif ext == ".json":
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
                pretty = json.dumps(data, indent=4)
                for line in pretty.split('\n'):
                    pdf.cell(200, 10, txt=line.strip(), ln=1)
            pdf.output(output_pdf)

        elif ext == ".xml":
            tree = ET.parse(filepath)
            root = tree.getroot()
            xml_str = ET.tostring(root, encoding='unicode')
            for line in xml_str.split('\n'):
                pdf.cell(200, 10, txt=line.strip(), ln=1)
            pdf.output(output_pdf)

        else:
            return "Unsupported file format"

        # ✅ Cleanup files after response
        @after_this_request
        def cleanup(response):
            def delete_files_later():
                time.sleep(5)
                for f in [filepath, output_pdf]:
                    try:
                        if os.path.exists(f):
                            os.remove(f)
                            print(f"Deleted: {f}")
                    except Exception as e:
                        print(f"Error deleting {f}: {e}")
            threading.Thread(target=delete_files_later).start()
            return response

        return send_file(output_pdf, as_attachment=True)

    except Exception as e:
        traceback.print_exc()
        return f"Error during conversion: {e}"

# ✅ Handle large file uploads
@app.errorhandler(RequestEntityTooLarge)
def handle_large_file(e):
    return "File too large. Maximum allowed size is 10 MB.", 413

if __name__ == '__main__':
    app.run(debug=True)
